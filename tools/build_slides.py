from __future__ import annotations

import csv
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/research/cbim/vast/sf895/code/Rutgers/cs550/FlowRanking")
ASSETS = ROOT / "assets"
FIG_DIR = ROOT / "static" / "images" / "movielens"
PAPER_DIR = Path("/research/cbim/vast/sf895/code/Rutgers/cs550/movielens_acm_paper")
DATA_DIR = Path("/research/cbim/vast/sf895/code/Rutgers/cs550/Class_Project/movielens_final_project/data/ml-latest-small")
OUT = ASSETS / "flow_ranking_slides.pptx"

TITLE_FONT = "Cambria"
BODY_FONT = "Aptos"

NAVY = RGBColor(17, 34, 68)
RED = RGBColor(187, 46, 55)
ROSE = RGBColor(246, 224, 222)
BLUE = RGBColor(229, 237, 248)
GOLD = RGBColor(179, 142, 74)
INK = RGBColor(35, 39, 45)
MUTED = RGBColor(98, 105, 117)
BG = RGBColor(247, 246, 243)
WHITE = RGBColor(255, 255, 255)
GRID = RGBColor(214, 217, 222)


def load_results():
    rows = {}
    with open(PAPER_DIR / "tables" / "results_summary.csv", newline="") as f:
        reader = csv.DictReader(f)
        for row in reader:
            rows[row["Model"]] = row
    return rows


def load_dataset_stats():
    ratings = pd.read_csv(DATA_DIR / "ratings.csv")
    movies = pd.read_csv(DATA_DIR / "movies.csv")
    return {
        "ratings": f"{len(ratings):,}",
        "users": f"{ratings.userId.nunique():,}",
        "movies": f"{ratings.movieId.nunique():,}",
        "mean_rating": f"{ratings.rating.mean():.2f}",
    }


def add_textbox(slide, left, top, width, height, text="", font_size=18, color=INK,
                bold=False, font_name=BODY_FONT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font_name
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, bullets, font_size=20, color=INK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    first = True
    for bullet in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = bullet
        p.level = 0
        p.bullet = True
        p.space_after = Pt(8)
        for run in p.runs:
            run.font.name = BODY_FONT
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
    return box


def add_header(slide, title, subtitle=""):
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.52))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    add_textbox(slide, Inches(0.45), Inches(0.07), Inches(7.8), Inches(0.26), title,
                font_size=26, color=WHITE, bold=True, font_name=TITLE_FONT)
    if subtitle:
        add_textbox(slide, Inches(8.65), Inches(0.1), Inches(4.2), Inches(0.2), subtitle,
                    font_size=11, color=WHITE, align=PP_ALIGN.RIGHT)


def add_footer(slide, page_num):
    add_textbox(
        slide,
        Inches(0.45),
        Inches(7.05),
        Inches(12.2),
        Inches(0.2),
        f"Flow Ranking | Sen Fang | CS550 Final Project | {page_num}",
        font_size=10,
        color=MUTED,
    )


def add_card(slide, left, top, width, height, title, lines, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = NAVY
    shape.line.width = Pt(1.0)
    add_textbox(slide, left + Inches(0.14), top + Inches(0.08), width - Inches(0.28), Inches(0.28),
                title, font_size=19, bold=True, font_name=TITLE_FONT, color=NAVY)
    add_bullets(slide, left + Inches(0.12), top + Inches(0.45), width - Inches(0.24), height - Inches(0.52),
                lines, font_size=14)


def add_results_table(slide, results):
    headers = ["Model", "MAE", "RMSE", "P@10", "R@10", "F@10", "NDCG@10"]
    models = ["ItemCF", "BiasedMF", "ClassicNeuMF", "FlowNeuMF"]
    table = slide.shapes.add_table(5, 7, Inches(0.42), Inches(0.95), Inches(12.45), Inches(2.55)).table
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.name = BODY_FONT
                r.font.bold = True
                r.font.size = Pt(13)
                r.font.color.rgb = WHITE
    for i, model in enumerate(models, start=1):
        row = [
            model,
            f"{float(results[model]['MAE']):.4f}",
            f"{float(results[model]['RMSE']):.4f}",
            f"{float(results[model]['Precision@10']):.4f}",
            f"{float(results[model]['Recall@10']):.4f}",
            f"{float(results[model]['F-measure@10']):.4f}",
            f"{float(results[model]['NDCG@10']):.4f}",
        ]
        for j, value in enumerate(row):
            cell = table.cell(i, j)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = ROSE if model == "FlowNeuMF" else WHITE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for r in p.runs:
                    r.font.name = BODY_FONT
                    r.font.size = Pt(12)
                    r.font.color.rgb = INK
                    r.font.bold = model == "FlowNeuMF"


def build():
    results = load_results()
    stats = load_dataset_stats()

    # Slide 1: Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    ribbon = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.28))
    ribbon.fill.solid()
    ribbon.fill.fore_color.rgb = NAVY
    ribbon.line.fill.background()
    add_textbox(slide, Inches(0.55), Inches(0.72), Inches(6.4), Inches(0.4), "Flow Ranking",
                font_size=30, bold=True, font_name=TITLE_FONT, color=RED)
    add_textbox(slide, Inches(0.55), Inches(1.18), Inches(6.1), Inches(0.9),
                "Flow-Regularized Neural Collaborative Filtering\nfor Movie Recommendation",
                font_size=24, bold=True, font_name=TITLE_FONT, color=NAVY)
    add_textbox(slide, Inches(0.58), Inches(2.3), Inches(4.4), Inches(1.1),
                "Sen Fang\nRutgers University CS550 Final Project\nsf895@scarletmail.rutgers.edu",
                font_size=17, color=INK)
    slide.shapes.add_picture(str(FIG_DIR / "topn_metrics_bar.png"), Inches(6.72), Inches(0.88), width=Inches(5.9))
    add_textbox(slide, Inches(6.75), Inches(6.45), Inches(5.8), Inches(0.22),
                "Final Top-10 comparison after fixing the evaluation mismatch.",
                font_size=11, color=MUTED, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.58), Inches(6.72), Inches(5.6), Inches(0.2),
                "Project page: fangsen9000.github.io/FlowRanking", font_size=11, color=MUTED)
    add_footer(slide, 1)

    # Slide 2: Setup
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    add_header(slide, "Problem Setup", "Why this project matters")
    add_bullets(slide, Inches(0.6), Inches(0.95), Inches(6.0), Inches(4.8), [
        "Goal: test whether flow-style regularization improves a neural recommender under a fair comparison.",
        "Dataset: MovieLens small with explicit ratings and per-user 80/20 holdout.",
        "Tasks: rating prediction and Top-10 recommendation.",
        "Metrics: MAE, RMSE, Precision@10, Recall@10, F-measure@10, and NDCG@10.",
    ], font_size=20)
    add_card(slide, Inches(7.15), Inches(1.05), Inches(2.5), Inches(1.45), "Ratings", [stats["ratings"]], BLUE)
    add_card(slide, Inches(9.85), Inches(1.05), Inches(2.5), Inches(1.45), "Users", [stats["users"]], BLUE)
    add_card(slide, Inches(7.15), Inches(2.75), Inches(2.5), Inches(1.45), "Movies", [stats["movies"]], BLUE)
    add_card(slide, Inches(9.85), Inches(2.75), Inches(2.5), Inches(1.45), "Avg Rating", [stats["mean_rating"]], BLUE)
    slide.shapes.add_picture(str(FIG_DIR / "rating_metrics_bar.png"), Inches(7.0), Inches(4.45), width=Inches(5.75))
    add_footer(slide, 2)

    # Slide 3: Methods
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    add_header(slide, "Compared Methods", "From classical CF to the proposed model")
    add_card(slide, Inches(0.5), Inches(1.0), Inches(3.0), Inches(2.0), "ItemCF", [
        "Item-item cosine similarity.",
        "Strong classical baseline for explicit ratings.",
    ], BLUE)
    add_card(slide, Inches(3.65), Inches(1.0), Inches(3.0), Inches(2.0), "BiasedMF", [
        "Matrix factorization with user/item bias terms.",
        "Strong latent-factor baseline.",
    ], BLUE)
    add_card(slide, Inches(6.8), Inches(1.0), Inches(3.0), Inches(2.0), "ClassicNeuMF", [
        "GMF branch + MLP tower.",
        "Neural baseline without flow regularization.",
    ], BLUE)
    add_card(slide, Inches(9.95), Inches(1.0), Inches(2.85), Inches(2.0), "FlowNeuMF", [
        "ClassicNeuMF backbone.",
        "Adds flow + consistency losses.",
    ], ROSE)
    add_bullets(slide, Inches(0.75), Inches(3.55), Inches(11.3), Inches(2.6), [
        "Neural models use the same embedding size, optimizer family, and train/test split.",
        "FlowNeuMF changes the regularization strategy rather than redefining the recommendation task.",
        "This makes the comparison against ClassicNeuMF the key test of the proposed idea.",
    ], font_size=18)
    add_footer(slide, 3)

    # Slide 4: Evaluation repair
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    add_header(slide, "Critical Debugging: Evaluation Repair", "Why the first ranking plots looked wrong")
    add_bullets(slide, Inches(0.6), Inches(0.98), Inches(6.1), Inches(5.0), [
        "Earlier Top-10 evaluation treated every held-out movie as relevant, including low-rated ones.",
        "That was inconsistent with the neural ranking loss, which used positive interactions only.",
        "The repaired protocol defines relevance as held-out ratings >= 4.0.",
        "ItemCF ranking was also repaired to use positive-history similarity propagation instead of pure rating-regression scores.",
    ], font_size=19)
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.15), Inches(1.15), Inches(5.1), Inches(2.3))
    panel.fill.solid()
    panel.fill.fore_color.rgb = WHITE
    panel.line.color.rgb = GOLD
    panel.line.width = Pt(1.2)
    add_textbox(slide, Inches(7.4), Inches(1.35), Inches(4.65), Inches(1.5),
                "Effect of the repair:\n\nItemCF Top-10 metrics moved from nearly zero to a reasonable competitive range, which fixed both the table and the bar charts.",
                font_size=18, color=INK)
    slide.shapes.add_picture(str(FIG_DIR / "precision_at_10.png"), Inches(7.05), Inches(3.95), width=Inches(5.35))
    add_footer(slide, 4)

    # Slide 5: Main results
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    add_header(slide, "Main Results", "Repaired full evaluation")
    add_results_table(slide, results)
    add_bullets(slide, Inches(0.6), Inches(3.9), Inches(5.9), Inches(2.2), [
        "ItemCF is best on MAE/RMSE.",
        "BiasedMF remains the strongest overall ranking baseline.",
        "FlowNeuMF beats ClassicNeuMF on every Top-10 metric.",
    ], font_size=19)
    slide.shapes.add_picture(str(FIG_DIR / "ndcg_at_10.png"), Inches(7.02), Inches(3.88), width=Inches(5.5))
    add_footer(slide, 5)

    # Slide 6: Ranking comparison
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    add_header(slide, "Ranking Quality Analysis", "Where the proposed method helps")
    slide.shapes.add_picture(str(FIG_DIR / "topn_metrics_bar.png"), Inches(0.52), Inches(1.0), width=Inches(7.1))
    add_bullets(slide, Inches(8.0), Inches(1.12), Inches(4.45), Inches(4.8), [
        "FlowNeuMF is not the best model overall, but it is a stronger neural recommender than ClassicNeuMF.",
        "The clearest gain is in NDCG@10, which is the strongest ranking-oriented evidence in the report.",
        "ItemCF still has slightly higher Recall@10, likely because neighborhood ranking is more popularity-driven.",
        "This makes the final claim narrower but more defensible.",
    ], font_size=18)
    add_footer(slide, 6)

    # Slide 7: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    add_header(slide, "Conclusion", "Submission-ready takeaway")
    add_bullets(slide, Inches(0.62), Inches(1.05), Inches(6.15), Inches(5.2), [
        "The evaluation pipeline is now consistent, reproducible, and visually interpretable.",
        "FlowNeuMF improves over ClassicNeuMF on all reported ranking metrics.",
        "The strongest baseline is still BiasedMF, so the contribution is an improved neural baseline rather than a new state of the art.",
        "Updated deliverables now align with each other: project page, ACM paper PDF, figures, source assets, and PPT.",
    ], font_size=20)
    slide.shapes.add_picture(str(FIG_DIR / "training_loss_curves.png"), Inches(7.0), Inches(1.22), width=Inches(5.55))
    add_textbox(slide, Inches(7.05), Inches(6.2), Inches(5.45), Inches(0.22),
                "Training curves kept for optimization discussion and ablation context.",
                font_size=11, color=MUTED, align=PP_ALIGN.CENTER)
    add_footer(slide, 7)

    prs.save(OUT)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


if __name__ == "__main__":
    build()
